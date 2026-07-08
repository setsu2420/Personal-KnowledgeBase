"""
MinerU 文档解析服务（Flask REST API）
内置于智能情报分析平台，提供高质量文档解析能力

支持格式：
- PDF → Markdown（MinerU pipeline / marker）
- HTML → Markdown（MinerU-HTML / BeautifulSoup fallback）
- DOCX/PPTX → Markdown（MinerU）
- 图片 → Markdown（MinerU OCR）

启动方式：
  python parse_document.py --port 8100

Java 后端通过 HTTP 调用此服务：
  POST /api/parse  - 解析本地文件
  POST /api/parse-url - 解析URL网页
  GET /api/health - 健康检查
"""

import argparse
import json
import os
import sys
import tempfile
import traceback
from pathlib import Path

from flask import Flask, request, jsonify
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

# 检测可用的解析器
import shutil as _shutil
_mineru_available = _shutil.which('mineru') is not None
_marker_available = False

if _mineru_available:
    print("[MinerU] mineru CLI found", file=sys.stderr)
else:
    # 检查 Python module
    try:
        import mineru
        _mineru_available = True
        print("[MinerU] mineru Python module available", file=sys.stderr)
    except ImportError:
        print("[MinerU] mineru NOT available", file=sys.stderr)

try:
    from marker.converters.pdf import PdfConverter
    _marker_available = True
    print("[MinerU] marker loaded successfully", file=sys.stderr)
except ImportError as e:
    print(f"[MinerU] marker not available: {e}", file=sys.stderr)

# PaddleOCR-VL 表格识别（本地 OCR 后备）
# 使用子进程安全检测 PaddleOCR 是否可用（NumPy 版本冲突可能导致 C 扩展崩溃）
_paddleocr_vl_available = False
_paddleocr_vl_pipe = None

import subprocess as _sp
_paddleocr_safe = False
try:
    _check = _sp.run(
        [sys.executable, '-c', 'import os; os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"]="True"; from paddleocr import PaddleOCR'],
        capture_output=True, timeout=15
    )
    _paddleocr_safe = (_check.returncode == 0)
except Exception:
    _paddleocr_safe = False

if _paddleocr_safe:
    try:
        import os
        os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")
        from paddleocr import PaddleOCR
        try:
            _test_ocr = PaddleOCR(use_textline_orientation=True, lang='ch')
        except (TypeError, ValueError):
            _test_ocr = PaddleOCR(use_angle_cls=True, lang='ch')
        _paddleocr_vl_available = True
        _paddleocr_vl_pipe = _test_ocr
        print("[PaddleOCR-VL] loaded successfully", file=sys.stderr)
    except Exception as e:
        print(f"[PaddleOCR-VL] init failed: {e}", file=sys.stderr)
else:
    print("[PaddleOCR-VL] NOT available (import check failed, likely NumPy version conflict)", file=sys.stderr)

# HTML 解析 fallback
try:
    from bs4 import BeautifulSoup
    import re
    _bs4_available = True
except ImportError:
    _bs4_available = False


def parse_pdf_mineru(file_path: str, **kwargs) -> dict:
    """使用 MinerU CLI 解析 PDF/DOCX/PPTX 文件"""
    import subprocess
    import shutil

    # 创建临时输出目录
    output_dir = tempfile.mkdtemp(prefix='mineru_out_')

    try:
        # 使用 MinerU CLI：mineru -p <input> -o <output> -b pipeline
        mineru_bin = shutil.which('mineru')
        if not mineru_bin:
            raise RuntimeError("mineru command not found in PATH")

        cmd = [
            mineru_bin,
            '-p', file_path,
            '-o', output_dir,
            '-b', 'pipeline',  # pipeline backend 不需要 GPU
        ]

        print(f"[MinerU CLI] Running: {' '.join(cmd)}", file=sys.stderr)
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=600,  # 10分钟超时
        )

        if result.returncode != 0:
            print(f"[MinerU CLI] returncode={result.returncode}", file=sys.stderr)
            print(f"[MinerU CLI] stderr: {result.stderr[:1000]}", file=sys.stderr)
            print(f"[MinerU CLI] stdout: {result.stdout[:500]}", file=sys.stderr)

        # 读取输出的 markdown 文件
        markdown = ""
        images = []

        output_path = Path(output_dir)
        # MinerU 在 output_dir 下创建以文件名命名的子目录
        for md_file in output_path.rglob('*.md'):
            markdown += md_file.read_text(encoding='utf-8') + '\n'

        # 收集图片信息
        for img_file in output_path.rglob('*.png'):
            rel_path = str(img_file.relative_to(output_path))
            images.append({'path': rel_path, 'caption': ''})
        for img_file in output_path.rglob('*.jpg'):
            rel_path = str(img_file.relative_to(output_path))
            images.append({'path': rel_path, 'caption': ''})

        if not markdown:
            # 尝试读取 content_list.json
            for json_file in output_path.rglob('content_list.json'):
                with open(json_file) as f:
                    content_list = json.load(f)
                    for item in content_list:
                        if item.get('type') == 'text':
                            markdown += item.get('text', '') + '\n'
                        elif item.get('type') == 'table':
                            markdown += item.get('html', '') + '\n'
                        elif item.get('type') == 'image':
                            images.append({
                                'path': item.get('img_path', ''),
                                'caption': item.get('img_caption', ''),
                            })

        return {
            'markdown': markdown,
            'images': images,
            'parser': 'mineru-cli',
        }
    finally:
        # 清理临时目录
        shutil.rmtree(output_dir, ignore_errors=True)


def parse_pdf_marker(file_path: str) -> dict:
    """使用 marker 解析 PDF 文件（备选方案）"""
    if not _marker_available:
        raise RuntimeError("marker not installed. Run: pip install marker-pdf")

    converter = PdfConverter()
    result = converter(file_path)

    return {
        'markdown': result.markdown,
        'images': [],
        'parser': 'marker',
    }


def parse_html_content(html_content: str) -> dict:
    """解析 HTML 内容为 Markdown"""
    if not _bs4_available:
        # 简单 fallback：去标签
        import re
        text = re.sub(r'<[^>]+>', ' ', html_content)
        text = re.sub(r'\s+', ' ', text).strip()
        return {'markdown': text, 'parser': 'fallback'}

    soup = BeautifulSoup(html_content, 'html.parser')

    # 移除 script/style
    for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
        tag.decompose()

    # 尝试提取主内容（article 或 main）
    main_content = soup.find('article') or soup.find('main') or soup.find('body') or soup
    if main_content is None:
        main_content = soup

    # 提取标题
    title = ''
    title_tag = soup.find('title')
    if title_tag:
        title = title_tag.get_text(strip=True)

    # 转换为 Markdown
    lines = []

    for element in main_content.descendants:
        if not hasattr(element, 'name'):
            continue

        if element.name == 'h1':
            lines.append(f"# {element.get_text(strip=True)}\n")
        elif element.name == 'h2':
            lines.append(f"## {element.get_text(strip=True)}\n")
        elif element.name == 'h3':
            lines.append(f"### {element.get_text(strip=True)}\n")
        elif element.name in ('h4', 'h5', 'h6'):
            level = int(element.name[1])
            lines.append(f"{'#' * level} {element.get_text(strip=True)}\n")
        elif element.name == 'p':
            text = element.get_text(strip=True)
            if text:
                lines.append(text + '\n')
        elif element.name == 'li':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"- {text}\n")
        elif element.name == 'blockquote':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"> {text}\n")
        elif element.name == 'code':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"`{text}`\n")
        elif element.name == 'pre':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"```\n{text}\n```\n")
        elif element.name == 'a':
            text = element.get_text(strip=True)
            href = element.get('href', '')
            if text and href:
                lines.append(f"[{text}]({href})\n")
        elif element.name == 'img':
            src = element.get('src', '')
            alt = element.get('alt', '')
            if src:
                lines.append(f"![{alt}]({src})\n")
        elif element.name == 'table':
            # 简单表格转 Markdown
            rows = element.find_all('tr')
            if rows:
                for i, row in enumerate(rows):
                    cells = row.find_all(['td', 'th'])
                    cell_texts = [c.get_text(strip=True) for c in cells]
                    lines.append('| ' + ' | '.join(cell_texts) + ' |')
                    if i == 0:
                        lines.append('| ' + ' | '.join(['---'] * len(cell_texts)) + ' |')
                lines.append('')
        elif element.name == 'br':
            lines.append('\n')

    # 去重和清理
    markdown = '\n'.join(lines)
    # 清理多余空行
    markdown = re.sub(r'\n{3,}', '\n\n', markdown)
    # 清理行内多余空格
    markdown = re.sub(r'  +', ' ', markdown)

    return {
        'markdown': markdown.strip(),
        'title': title,
        'parser': 'beautifulsoup',
    }


def parse_html_mineru(html_path: str) -> dict:
    """尝试使用 MinerU 解析 HTML"""
    try:
        converter = DocumentConverter()
        result = converter(html_path)
        markdown = ""
        if hasattr(result, 'md_content'):
            markdown = result.md_content
        elif hasattr(result, 'content_list'):
            for item in result.content_list:
                if item.get('type') == 'text':
                    markdown += item.get('text', '') + '\n'
        return {
            'markdown': markdown,
            'parser': 'mineru-html',
        }
    except Exception:
        # fallback to BeautifulSoup
        with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
            return parse_html_content(f.read())


@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'ok',
        'mineru': _mineru_available,
        'marker': _marker_available,
        'bs4': _bs4_available,
        'paddleocr_vl': _paddleocr_vl_available,
    })


def _ocr_result_to_markdown(ocr_result) -> str:
    """将 PaddleOCR 识别结果转换为 Markdown 表格"""
    if not ocr_result:
        return ""

    # PaddleOCR 返回格式: [[{box, text, confidence}, ...], ...]
    # 每个子列表是一行文字的识别结果
    lines = []
    for line_group in ocr_result:
        if not isinstance(line_group, list):
            continue
        # 按 x 坐标排序，提取同一行的文本
        boxes = []
        for item in line_group:
            if isinstance(item, dict) and 'text' in item:
                box = item.get('box', [[0, 0]])
                y_pos = box[0][1] if box else 0
                x_pos = box[0][0] if box else 0
                boxes.append((x_pos, y_pos, item['text']))
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                # 旧版格式: [[[x,y],...], (text, confidence)]
                box = item[0]
                text_conf = item[1]
                text = text_conf[0] if isinstance(text_conf, (list, tuple)) else str(text_conf)
                x_pos = box[0][0] if box and len(box) > 0 else 0
                y_pos = box[0][1] if box and len(box) > 0 else 0
                boxes.append((x_pos, y_pos, text))

        # 按 y 坐标分行，x 坐标排序
        boxes.sort(key=lambda b: (round(b[1] / 20) * 20, b[0]))

        # 按行分组（y 坐标接近的归为一行）
        rows = []
        current_row = []
        last_y = -1
        for x, y, text in boxes:
            if last_y >= 0 and abs(y - last_y) > 15:
                if current_row:
                    rows.append(current_row)
                current_row = []
            current_row.append((x, text))
            last_y = y
        if current_row:
            rows.append(current_row)

        # 转换为 Markdown 表格
        for i, row in enumerate(rows):
            row.sort(key=lambda c: c[0])
            cells = [c[1].strip().replace('|', '\\|') for c in row]
            lines.append('| ' + ' | '.join(cells) + ' |')
            if i == 0:
                lines.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')

    return '\n'.join(lines)


@app.route('/api/ocr-table', methods=['POST'])
def ocr_table():
    """
    使用 PaddleOCR-VL 识别表格图片为 Markdown 表格

    接收：
    - multipart: image 字段（图片文件）
    - 或 JSON: {"image": "base64..."}

    返回：
    - markdown: Markdown 格式的表格
    """
    if not _paddleocr_vl_available:
        return jsonify({
            'error': 'PaddleOCR-VL 未安装。请运行: pip install paddleocr 或 pip install transformers && pip install paddlepaddle',
            'install_hint': 'pip install paddleocr'
        }), 503

    image_data = None

    # 尝试 multipart 上传
    if 'image' in request.files:
        image_data = request.files['image'].read()
    # 尝试 JSON base64
    elif request.is_json:
        data = request.get_json()
        if 'image' in data:
            import base64
            image_data = base64.b64decode(data['image'])

    if not image_data:
        return jsonify({'error': '请提供图片文件（multipart image 字段）或 base64 JSON'}), 400

    # 保存到临时文件
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(image_data)
        tmp_path = tmp.name

    try:
        # 使用 PaddleOCR 识别
        result = _paddleocr_vl_pipe.ocr(tmp_path, cls=True)

        if not result:
            return jsonify({'error': 'OCR 未识别到任何文本', 'markdown': ''})

        markdown = _ocr_result_to_markdown(result)

        return jsonify({
            'markdown': markdown,
            'parser': 'paddleocr-vl',
            'line_count': len(markdown.split('\n')) if markdown else 0
        })

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500
    finally:
        import os
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


@app.route('/api/parse', methods=['POST'])
def parse_file():
    """
    解析本地文件
    
    请求参数：
    - file_path: 文件路径
    - parser: 解析器选择 (mineru/marker/auto)，默认 auto
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - images: 提取的图片列表
    - parser: 使用的解析器
    - title: 文档标题（HTML 文档）
    """
    data = request.get_json()
    if not data or 'file_path' not in data:
        return jsonify({'error': 'file_path is required'}), 400

    file_path = data['file_path']
    parser_choice = data.get('parser', 'auto')

    if not os.path.exists(file_path):
        return jsonify({'error': f'File not found: {file_path}'}), 404

    ext = Path(file_path).suffix.lower()

    try:
        if ext == '.pdf':
            if parser_choice == 'marker' and _marker_available:
                result = parse_pdf_marker(file_path)
            elif _mineru_available:
                result = parse_pdf_mineru(file_path)
            elif _marker_available:
                result = parse_pdf_marker(file_path)
            else:
                # 最终 fallback：尝试 PyPDF2 或 pdfminer
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    text = '\n'.join([page.extract_text() or '' for page in reader.pages])
                    result = {'markdown': text, 'parser': 'pypdf'}
                except ImportError:
                    raise RuntimeError("No PDF parser available. Install mineru or marker-pdf or pypdf")

        elif ext in ('.html', '.htm'):
            if parser_choice == 'mineru' and _mineru_available:
                result = parse_html_mineru(file_path)
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result = parse_html_content(f.read())

        elif ext in ('.docx', '.doc'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    from docx import Document as DocxDocument
                    doc = DocxDocument(file_path)
                    text = '\n'.join([p.text for p in doc.paragraphs])
                    result = {'markdown': text, 'parser': 'python-docx'}
                except ImportError:
                    raise RuntimeError("No DOCX parser available. Install mineru or python-docx")

        elif ext in ('.pptx', '.ppt'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                text_parts.append(shape.text)
                    result = {'markdown': '\n'.join(text_parts), 'parser': 'python-pptx'}
                except ImportError:
                    raise RuntimeError("No PPTX parser available. Install mineru or python-pptx")

        elif ext in ('.xlsx', '.xls'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    md_parts = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        md_parts.append(f"## {sheet}\n")
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else '' for c in row]
                            md_parts.append('| ' + ' | '.join(cells) + ' |')
                        md_parts.append('')
                    result = {'markdown': '\n'.join(md_parts), 'parser': 'openpyxl'}
                except ImportError:
                    raise RuntimeError("No XLSX parser available. Install mineru or openpyxl")

        elif ext in ('.txt', '.md', '.markdown', '.csv'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                result = {'markdown': f.read(), 'parser': 'plaintext'}

        else:
            return jsonify({'error': f'Unsupported file type: {ext}'}), 400

        return jsonify(result)

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e), 'traceback': traceback.format_exc()}), 500


@app.route('/api/parse-url', methods=['POST'])
def parse_url():
    """
    解析URL网页内容
    
    请求参数：
    - url: 网页URL
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - title: 网页标题
    - parser: 使用的解析器
    """
    data = request.get_json()
    if not data or 'url' not in data:
        return jsonify({'error': 'url is required'}), 400

    url = data['url']

    try:
        import urllib.request
        req = urllib.request.Request(
            url,
            headers={'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'}
        )
        with urllib.request.urlopen(req, timeout=30) as response:
            html_content = response.read().decode('utf-8', errors='ignore')

        result = parse_html_content(html_content)
        result['url'] = url
        return jsonify(result)

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/parse-html', methods=['POST'])
def parse_html_direct():
    """
    直接解析 HTML 内容（不需要文件）
    
    请求参数：
    - html: HTML 字符串
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - title: 网页标题
    """
    data = request.get_json()
    if not data or 'html' not in data:
        return jsonify({'error': 'html content is required'}), 400

    try:
        result = parse_html_content(data['html'])
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='MinerU Document Parser Service')
    parser.add_argument('--port', type=int, default=8100, help='Port to listen on')
    parser.add_argument('--host', type=str, default='127.0.0.1', help='Host to bind to')
    args = parser.parse_args()

    print(f"[MinerU Service] Starting on {args.host}:{args.port}")
    print(f"  MinerU: {'available' if _mineru_available else 'NOT available'}")
    print(f"  marker: {'available' if _marker_available else 'NOT available'}")
    print(f"  bs4: {'available' if _bs4_available else 'NOT available'}")
    print(f"  PaddleOCR-VL: {'available' if _paddleocr_vl_available else 'NOT available'}")

    app.run(host=args.host, port=args.port, debug=False)
